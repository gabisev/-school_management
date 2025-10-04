#!/usr/bin/env python3
"""
Script pour créer une présentation PowerPoint du système de gestion scolaire
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    """Crée une présentation PowerPoint complète"""
    
    # Créer une nouvelle présentation
    prs = Presentation()
    
    # Définir les couleurs du thème
    primary_color = RGBColor(52, 152, 219)  # Bleu
    secondary_color = RGBColor(44, 62, 80)  # Gris foncé
    accent_color = RGBColor(46, 204, 113)   # Vert
    
    # Slide 1: Page de titre
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "🎓 Système de Gestion Scolaire"
    subtitle.text = "Solution Complète pour la Gestion Éducative\nDéveloppé avec Django - Interface Moderne - Multi-utilisateurs"
    
    # Slide 2: Vue d'ensemble
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "📋 Vue d'Ensemble du Projet"
    content2.text = """🎯 Objectif
Créer une solution complète de gestion scolaire qui simplifie l'administration éducative et améliore la communication entre tous les acteurs de l'école.

👥 Utilisateurs Cibles
• 👨‍🎓 Élèves - Consultation des notes, bulletins, emploi du temps
• 👨‍🏫 Professeurs - Saisie des notes, gestion des classes
• 👨‍👩‍👧‍👦 Parents - Suivi de la scolarité de leurs enfants
• 👨‍💼 Administrateurs - Gestion complète du système"""
    
    # Slide 3: Fonctionnalités principales
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "⚡ Fonctionnalités Principales"
    content3.text = """📚 Gestion Académique
• Gestion des matières et programmes
• Système de notes et calcul automatique
• Génération de bulletins trimestriels
• Gestion des évaluations et examens

👥 Gestion des Utilisateurs
• Profils personnalisés
• Authentification sécurisée multi-niveaux
• Permissions granulaires par rôle
• Organisation des élèves par classe

📅 Planning et Organisation
• Emploi du temps interactif
• Réservation de salles
• Calendrier des événements scolaires"""
    
    # Slide 4: Architecture technique
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "🏗️ Architecture Technique"
    content4.text = """🔧 Stack Technologique
• Frontend: HTML5, CSS3, JavaScript, Bootstrap 4
• Backend: Django 4.2.5, Python 3.8+
• Base de données: SQLite (dev), PostgreSQL/MySQL (prod)
• Déploiement: Docker, Nginx, Gunicorn

📊 Caractéristiques Techniques
• Architecture MVC - Séparation claire des responsabilités
• API REST - Interface programmatique complète
• Responsive Design - Compatible mobile et desktop
• Sécurité - Authentification, autorisation, audit"""
    
    # Slide 5: Interface utilisateur
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "🎨 Interface Utilisateur"
    content5.text = """📱 Design Responsive
• Desktop - Interface complète avec toutes les fonctionnalités
• Tablette - Optimisé pour la navigation tactile
• Mobile - Version simplifiée pour l'accès rapide

🎯 Tableaux de Bord Personnalisés
• Élève: Notes récentes, prochains cours, messages
• Professeur: Classes assignées, notes à saisir, planning
• Parent: Enfants, notes, absences, communications
• Administrateur: Statistiques, gestion utilisateurs

💬 Système de Communication
• Messagerie intégrée entre utilisateurs
• Notifications en temps réel
• Annonces et communications de masse"""
    
    # Slide 6: Déploiement et DevOps
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    content6 = slide6.placeholders[1]
    
    title6.text = "🚀 Déploiement et DevOps"
    content6.text = """🐳 Containerisation Docker
• Démarrage rapide: docker-compose up -d
• Environnements multiples (dev, staging, prod)
• Configuration Nginx pour la production

🔄 CI/CD avec GitHub Actions
• Tests automatisés à chaque commit
• Qualité du code (linting, formatting)
• Scan de sécurité des vulnérabilités
• Déploiement automatique

📊 Monitoring et Sécurité
• Logs centralisés et métriques de performance
• HTTPS et authentification sécurisée
• Audit et traçabilité des actions"""
    
    # Slide 7: Installation et utilisation
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    
    title7.text = "⚙️ Installation et Utilisation"
    content7.text = """🚀 Installation Rapide
1. git clone https://github.com/gabisev/-school_management.git
2. python setup_dev.py
3. python manage.py runserver

🐳 Avec Docker (Recommandé)
• Développement: docker-compose -f docker-compose.dev.yml up -d
• Production: docker-compose -f docker-compose.prod.yml up -d

🎯 Premiers Pas
1. Accédez à http://localhost:8000
2. Connectez-vous avec admin/admin123
3. Créez vos premières classes et matières
4. Ajoutez des utilisateurs
5. Configurez l'emploi du temps"""
    
    # Slide 8: Avantages et bénéfices
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    content8 = slide8.placeholders[1]
    
    title8.text = "💡 Avantages et Bénéfices"
    content8.text = """🎯 Pour l'École
• Efficacité - Automatisation des tâches administratives
• Économies - Réduction des coûts de gestion
• Traçabilité - Historique complet des actions
• Sécurité - Protection des données sensibles

👨‍🏫 Pour les Professeurs
• Gain de temps - Saisie simplifiée des notes
• Communication - Contact direct avec parents
• Organisation - Planning et gestion des classes
• Suivi - Évolution des élèves en temps réel

👨‍👩‍👧‍👦 Pour les Parents
• Transparence - Accès direct aux informations
• Communication - Contact facilité avec l'école
• Suivi - Monitoring de la progression
• Convenance - Accès 24/7 depuis n'importe où"""
    
    # Slide 9: Roadmap
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    content9 = slide9.placeholders[1]
    
    title9.text = "🗺️ Roadmap et Évolutions"
    content9.text = """📅 Q1 2024 - Améliorations Immédiates
• API REST complète
• Interface mobile native
• Notifications push
• Thème sombre

🚀 Q2-Q3 2024 - Nouvelles Fonctionnalités
• Module bibliothèque
• Système de paiement des frais
• Transport scolaire
• Intelligence artificielle

🌟 Q4 2024+ - Vision Long Terme
• Multi-tenant (plusieurs établissements)
• Internationalisation
• Analytics avancés
• Intégrations externes"""
    
    # Slide 10: Conclusion
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    content10 = slide10.placeholders[1]
    
    title10.text = "🎯 Conclusion"
    content10.text = """✅ Ce que nous avons accompli
• 185 fichiers organisés professionnellement
• 15+ modules fonctionnels
• 4 environnements de déploiement
• 100% Open Source

🎊 Points Forts du Projet
• Solution complète pour tous les aspects scolaires
• Technologies modernes et à jour
• Documentation complète et détaillée
• Déploiement facile avec Docker
• Architecture extensible et évolutive

📞 Contact et Ressources
• Dépôt: github.com/gabisev/-school_management
• Documentation: README.md, QUICK_START.md
• Support: Issues GitHub, Discussions
• Licence: MIT - Libre d'utilisation

🚀 Prêt à révolutionner la gestion scolaire !"""
    
    # Sauvegarder la présentation
    filename = "Système_Gestion_Scolaire_Presentation.pptx"
    prs.save(filename)
    print(f"✅ Présentation créée: {filename}")
    
    return filename

if __name__ == "__main__":
    try:
        # Vérifier si python-pptx est installé
        import pptx
        create_presentation()
    except ImportError:
        print("❌ Le module python-pptx n'est pas installé.")
        print("📦 Installez-le avec: pip install python-pptx")
        print("🔄 Ou utilisez le fichier HTML: presentation.html")
